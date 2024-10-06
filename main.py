import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import nltk
from nltk.sentiment import SentimentIntensityAnalyzer
from sklearn.feature_extraction.text import CountVectorizer

# Download necessary NLTK resources
nltk.download('vader_lexicon')

def analyze_notes(notes):
    sia = SentimentIntensityAnalyzer()
    sentiments = []
    for note in notes:
        # Ensure the note is a string before analysis
        if isinstance(note, str):
            sentiment_score = sia.polarity_scores(note)
            sentiments.append((note, sentiment_score))
        else:
            # Skip non-string notes or handle them accordingly
            sentiments.append((note, {"compound": 0}))
    
    return sentiments

def extract_keywords(notes, top_n=5):
    # Filter out non-string notes
    string_notes = [note for note in notes if isinstance(note, str)]
    
    if len(string_notes) == 0:
        return []

    vectorizer = CountVectorizer(stop_words='english', max_features=top_n)
    X = vectorizer.fit_transform(string_notes)
    keywords = vectorizer.get_feature_names_out()
    return keywords

def generate_observations(notes):
    # Analyze Sentiment
    sentiments = analyze_notes(notes)
    
    positive_notes = [note for note, score in sentiments if score['compound'] > 0]
    negative_notes = [note for note, score in sentiments if score['compound'] < 0]
    neutral_notes = [note for note, score in sentiments if score['compound'] == 0]

    # Extract Key Topics
    keywords = extract_keywords(notes)
    
    # Create observations summary
    summary = f"Observations Summary:\n\n"
    summary += f"- Total Notes: {len(notes)}\n"
    summary += f"- Positive Notes: {len(positive_notes)}\n"
    summary += f"- Negative Notes: {len(negative_notes)}\n"
    summary += f"- Neutral Notes: {len(neutral_notes)}\n"
    summary += f"- Key Topics: {', '.join(keywords)}\n"
    
    return summary

def excel_to_powerpoint(excel_file, ppt_template=None):
    # Load the Excel file
    wb = openpyxl.load_workbook(excel_file, data_only=True)
    sheet = wb.active
    
    # Create or Load PowerPoint Presentation
    if ppt_template:
        prs = Presentation(ppt_template)
    else:
        prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Weekly Financial Report"
    subtitle.text = f"Data from: {excel_file}"
    
    # Extract financial data (costs in Column F, notes in Column E)
    notes = [row[4] for row in sheet.iter_rows(min_row=2, max_row=sheet.max_row, values_only=True) if row[4]]
    income_data = [row for row in sheet.iter_rows(min_row=2, max_row=sheet.max_row, values_only=True) if row[6] == "Income"]
    expense_data = [row for row in sheet.iter_rows(min_row=2, max_row=sheet.max_row, values_only=True) if row[6] == "Expense"]
    
    total_income = sum(row[5] for row in income_data)
    total_expenses = sum(row[5] for row in expense_data)
    net_profit = total_income - total_expenses
    
    # Add financial summary slide
    summary = f"Total Income: ${total_income:,.2f}\n"
    summary += f"Total Expenses: ${total_expenses:,.2f}\n"
    summary += f"Net Profit: ${net_profit:,.2f}"
    
    def add_text_slide(title, content):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = title
        tf = body_shape.text_frame
        tf.text = content

    add_text_slide("Financial Summary", summary)
    
    # Add revenue vs expenses chart slide
    chart_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(chart_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Revenue vs Expenses'
    
    # Add chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Revenue', 'Expenses']
    chart_data.add_series('Amount', (total_income, total_expenses))
    
    # Add chart to slide
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # Generate AI-powered analysis of notes
    observations = generate_observations(notes)
    
    # Add AI-generated observations slide
    add_text_slide("AI-Generated Observations", observations)
    
    # Save the presentation
    ppt_file = excel_file.replace('.xlsx', '_analysis.pptx')
    prs.save(ppt_file)
    print(f"PowerPoint presentation saved as {ppt_file}")

# Usage
excel_to_powerpoint('weekly_financial_report.xlsx')
